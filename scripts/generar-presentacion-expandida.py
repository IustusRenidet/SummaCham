# -*- coding: utf-8 -*-
"""
Genera la presentación técnica expandida de SummaCham.
~55 slides con información ultra-detallada.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# === COLORES ===
AZUL_OSCURO = RGBColor(0x0d, 0x47, 0xa1)
AZUL_MEDIO = RGBColor(0x1e, 0x88, 0xe5)
AZUL_CLARO = RGBColor(0x60, 0xa5, 0xfa)
GRIS_OSCURO = RGBColor(0x33, 0x33, 0x33)
GRIS_MEDIO = RGBColor(0x66, 0x66, 0x66)
GRIS_CLARO = RGBColor(0x94, 0xa3, 0xb8)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
VERDE = RGBColor(0x22, 0xc5, 0x5e)
NARANJA = RGBColor(0xFF, 0x98, 0x00)
ROJO = RGBColor(0xEF, 0x53, 0x50)
NEGRO = RGBColor(0x00, 0x00, 0x00)
FONDO_TITULO = RGBColor(0x0d, 0x47, 0xa1)
FONDO_SECCION = RGBColor(0x1a, 0x23, 0x7e)

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(6.25)

# === HELPERS ===
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=GRIS_OSCURO, alignment=PP_ALIGN.LEFT,
                font_name='Segoe UI'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=11,
                  color=GRIS_OSCURO, line_spacing=1.15, font_name='Segoe UI',
                  bold_first=False):
    """lines = list of strings"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(2)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def title_slide(title, subtitle="", bg_color=FONDO_TITULO):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, bg_color)
    add_textbox(slide, 0.5, 1.8, 9, 1.2, title, 36, True, BLANCO, PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 0.5, 3.2, 9, 1, subtitle, 16, False, AZUL_CLARO, PP_ALIGN.CENTER)
    return slide

def section_slide(part_num, title, bg_color=FONDO_SECCION):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, bg_color)
    add_textbox(slide, 0.5, 2.0, 9, 0.8, f"PARTE {part_num}", 14, False, AZUL_CLARO, PP_ALIGN.CENTER)
    add_textbox(slide, 0.5, 2.6, 9, 1.2, title, 30, True, BLANCO, PP_ALIGN.CENTER)
    return slide

def content_slide(title, body_lines, title_size=22, body_size=11):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.4, 0.2, 9.2, 0.7, title, title_size, True, AZUL_OSCURO, PP_ALIGN.LEFT)
    # Línea separadora
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1, Inches(0.4), Inches(0.85), Inches(9.2), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_CLARO
    shape.line.fill.background()
    add_multiline(slide, 0.4, 1.0, 9.2, 5.0, body_lines, body_size)
    return slide

def two_col_slide(title, left_lines, right_lines, title_size=22, body_size=11):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_textbox(slide, 0.4, 0.2, 9.2, 0.7, title, title_size, True, AZUL_OSCURO)
    shape = slide.shapes.add_shape(1, Inches(0.4), Inches(0.85), Inches(9.2), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZUL_CLARO
    shape.line.fill.background()
    add_multiline(slide, 0.4, 1.0, 4.3, 5.0, left_lines, body_size)
    add_multiline(slide, 5.0, 1.0, 4.6, 5.0, right_lines, body_size)
    return slide


# ============================================================
# SLIDE 1: PORTADA
# ============================================================
s = title_slide("SummaCham - PanelAMCHAM", "Documentación Técnica Completa y Detallada")
add_textbox(s, 0.5, 4.2, 9, 0.5, "Arquitectura, APIs, Queries, Plantillas, Gráficas, Base de Datos y más",
            13, False, AZUL_CLARO, PP_ALIGN.CENTER)
add_textbox(s, 0.5, 5.2, 9, 0.5, "Versión 4.1.0  |  Febrero 2026  |  Iustus Renidet",
            11, False, GRIS_CLARO, PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: QUE ES
# ============================================================
content_slide("¿Qué es PanelAmcham (SummaCham)?", [
    "Es una aplicación de escritorio para gestionar presupuestos financieros.",
    "",
    "Imagina un Excel inteligente que:",
    "  • Se conecta al sistema contable (Aspel COI / Firebird) para leer saldos reales",
    "  • Permite que varios usuarios trabajen al mismo tiempo con control de permisos",
    "  • Tiene un flujo de aprobación: cargar → revisar → aprobar → guardar en COI",
    "  • Genera reportes consolidados (SUMMARY y RESUMEN) automáticamente",
    "  • Crea gráficas interactivas (barras, líneas, pie, dona, radar, área)",
    "  • Exporta a Excel y PDF con un clic",
    "",
    "USUARIOS TÍPICOS:",
    "  • Analistas: Cargan datos de presupuesto (RH, Finanzas, Eventos, Comités, etc.)",
    "  • Supervisores: Revisan que los datos estén correctos",
    "  • Directores: Aprueban los presupuestos finales",
    "  • Administradores: Configuran usuarios, permisos, conexiones y plantillas",
    "",
    "CÓMO SE ACCEDE:",
    "  • Local: http://localhost:3005 (en la misma computadora, vía Electron)",
    "  • Remoto: https://panelamcham.amcham.com.mx (desde cualquier lugar, vía navegador)",
])

# ============================================================
# SLIDE 3: CONTENIDO
# ============================================================
two_col_slide("Contenido de esta Presentación", [
    "PARTE 1:  Arquitectura General",
    "PARTE 2:  Proceso Electron (Desktop)",
    "PARTE 3:  Express: Middleware y Servidor",
    "PARTE 4:  Login, Cookies, Sesiones",
    "PARTE 5:  Sistema de Permisos (3 niveles)",
    "PARTE 6:  ¿Qué es una API? + Ejemplo real",
    "PARTE 7:  APIs de Autenticación",
    "PARTE 8:  APIs de Reportes",
    "PARTE 9:  APIs de Presupuestos",
    "PARTE 10: APIs de Borradores (Drafts)",
], [
    "PARTE 11: APIs de Usuarios y Perfil",
    "PARTE 12: APIs de Plantillas (Layouts)",
    "PARTE 13: APIs de Saldos, Módulos y más",
    "PARTE 14: Mapa completo de TODAS las APIs",
    "PARTE 15: Firebird: De dónde salen los números",
    "PARTE 16: Cálculos: Real, Presupuesto, YTD",
    "PARTE 17: Motor de Reportes (Engine)",
    "PARTE 18: Gestor de Plantillas (14,700+ líneas)",
    "PARTE 19: Sistema de Gráficas",
    "PARTE 20: Base de Datos, Backups, Seguridad",
])

# ============================================================
# PARTE 1: ARQUITECTURA
# ============================================================
section_slide(1, "Arquitectura General")

content_slide("Cómo está construido el sistema", [
    "El sistema tiene 3 grandes partes que se comunican entre sí:",
    "",
    "┌──────────────────┐    ┌──────────────────┐    ┌──────────────────┐",
    "│   FRONTEND        │    │   BACKEND         │    │   BASES DE DATOS  │",
    "│ (Lo que VE el     │◄──►│ (El servidor que  │◄──►│                   │",
    "│  usuario)         │    │  procesa)         │    │ SQLite + Firebird │",
    "│ HTML+JS+CSS       │    │ Node.js+Express   │    │                   │",
    "│ Electron shell    │    │ Puerto 3005       │    │                   │",
    "└──────────────────┘    └──────────────────┘    └──────────────────┘",
    "",
    "PASO A PASO:",
    "1. El usuario abre la app → Electron abre ventana con navegador integrado",
    "2. El navegador carga http://localhost:3005",
    "3. Cuando necesita datos, el frontend envía una PETICIÓN HTTP al backend",
    "   Ejemplo: 'Dame el reporte RESUMEN de empresa1 para Junio 2025'",
    "4. El backend recibe la petición y busca datos en 2 lugares:",
    "   • SQLite: Qué cuentas mostrar, qué fórmulas aplicar, permisos del usuario",
    "   • Firebird: Los saldos reales del sistema contable Aspel COI",
    "5. El backend responde con los datos en formato JSON",
    "6. El frontend muestra los datos en tablas y gráficas",
])

two_col_slide("Tecnologías utilizadas", [
    "FRONTEND (lo que ve el usuario):",
    "",
    "  Electron 39.2.7  →  App de escritorio",
    "  Bootstrap 5       →  Diseño profesional",
    "  Chart.js 4.3      →  Gráficas interactivas",
    "  DataTables 1.13   →  Tablas con filtros",
    "  SweetAlert2        →  Diálogos bonitos",
    "  ExcelJS / XLSX    →  Exportar a Excel",
    "  jsPDF             →  Exportar a PDF",
    "  Font Awesome 6.4  →  Íconos",
    "",
    "Nota: El frontend usa HTML/JavaScript",
    "puro (sin React ni Angular).",
], [
    "BACKEND (el servidor):",
    "",
    "  Node.js            →  Motor JavaScript",
    "  Express.js         →  Framework web",
    "  Better-SQLite3     →  BD local (síncrona)",
    "  node-firebird      →  Conexión a Aspel COI",
    "  bcryptjs           →  Encriptar contraseñas",
    "  jsonwebtoken       →  Tokens JWT",
    "  Joi                →  Validar datos entrada",
    "  Helmet             →  Seguridad HTTP",
    "  express-session    →  Sesiones de usuario",
    "  express-rate-limit →  Limitar intentos",
    "  compression        →  Compresión gzip",
    "  cors               →  Control de orígenes",
])

# ============================================================
# PARTE 2: ELECTRON
# ============================================================
section_slide(2, "Proceso Electron (Desktop)")

content_slide("Cómo funciona Electron", [
    "Electron convierte una página web en una aplicación de escritorio nativa.",
    "SummaCham usa Electron 39.2.7 para correr en Windows.",
    "",
    "PROCESO PRINCIPAL (main.js):",
    "  1. Crea una ventana BrowserWindow con navegador Chromium integrado",
    "  2. Inicia el servidor Express en puerto 3005 (dentro del mismo proceso)",
    "  3. Carga http://localhost:3005 en la ventana",
    "  4. Configura auto-update para verificar nuevas versiones en GitHub Releases",
    "",
    "CONFIGURACIÓN DE LA VENTANA:",
    "  • Tamaño: 1280x800 mínimo, maximizado por defecto",
    "  • nodeIntegration: false (seguridad - JS del frontend NO puede acceder a Node)",
    "  • contextIsolation: true (aislamiento entre procesos)",
    "  • webSecurity: true (protección contra XSS)",
    "",
    "EMPAQUETADO (electron-builder):",
    "  • Target: NSIS installer (SummaCham Setup 4.1.0.exe)",
    "  • También genera portable (.exe sin instalación)",
    "  • Auto-update vía GitHub Releases (verifica, descarga, instala al cerrar)",
    "  • Los DATOS del usuario se conservan en userData/datos/ (solo se actualiza el código)",
])

# ============================================================
# PARTE 3: EXPRESS
# ============================================================
section_slide(3, "Express: Middleware y Servidor")

content_slide("Cadena de Middleware (orden exacto)", [
    "Cada petición HTTP pasa por esta cadena de middleware en src/server.js:",
    "",
    "  1. helmet()          → Cabeceras de seguridad (X-XSS-Protection, HSTS, etc.)",
    "  2. compression()     → Compresión gzip para respuestas grandes",
    "  3. cors({             → Control de orígenes permitidos",
    "       origin: [...],     Permite localhost:3005 y panelamcham.amcham.com.mx",
    "       credentials: true  Permite enviar cookies cross-origin",
    "     })",
    "  4. express.json({limit:'20mb'})   → Parsear cuerpo JSON (hasta 20MB)",
    "  5. express.urlencoded()            → Parsear formularios",
    "  6. express-session({   → Manejo de sesiones",
    "       store: SQLiteStore  Guarda sesiones en SQLite (no memoria)",
    "       name: 'panelamcham.sid'",
    "       cookie: { httpOnly, secure, maxAge:30min, rolling:true }",
    "     })",
    "  7. express-rate-limit   → Login: máx 10 intentos/10min, bloqueo 5min",
    "  8. express.static()     → Archivos estáticos (HTML, CSS, JS, imágenes)",
    "  9. Rutas API (/api/*)   → Cada ruta tiene su propio requireAuth + Joi",
])

content_slide("Registro de Rutas en el Servidor", [
    "Todas las rutas se registran en src/server.js así:",
    "",
    "  app.use('/api/auth',            require('./routes/auth'))           // Login/Logout",
    "  app.use('/api/usuarios',        require('./routes/usuarios'))       // CRUD usuarios",
    "  app.use('/api/perfil',          require('./routes/perfil'))         // Mi perfil",
    "  app.use('/api/empresas',        require('./routes/empresas'))       // Lista empresas",
    "  app.use('/api/modulos',         require('./routes/modulos'))        // Info módulos",
    "  app.use('/api/reportes',        require('./routes/reportes'))       // Summary/Resumen",
    "  app.use('/api/presupuestos',    require('./routes/presupuestos'))   // Presupuestos",
    "  app.use('/api/borradores',      require('./routes/borradores'))     // Drafts",
    "  app.use('/api/saldos',          require('./routes/saldos'))         // Saldos Firebird",
    "  app.use('/api/layouts-config',  require('./routes/layoutRoutes'))   // Plantillas",
    "  app.use('/api/insercion',       require('./routes/insercion'))      // Inserción filas",
    "  app.use('/api/graficas-config', require('./routes/graficas-config'))// Gráficas",
    "  app.use('/api/backups',         require('./routes/backups'))        // Backups",
    "  app.use('/api/notificaciones',  require('./routes/notificaciones')) // Alertas",
    "  app.use('/api/firebird-config', require('./routes/firebird-config'))// Config Firebird",
    "  app.use('/api/comentarios',     require('./routes/comentarios'))    // Comentarios",
    "",
    "  TOTAL: 16+ grupos de rutas, 80+ endpoints individuales",
])

# ============================================================
# PARTE 4: LOGIN
# ============================================================
section_slide(4, "Login, Cookies y Sesiones")

content_slide("Cómo funciona el Login (paso a paso)", [
    "Cuando escribes tu usuario y contraseña, pasa esto:",
    "",
    "  1. Frontend envía:  POST /api/auth/login",
    "     Body: { usuario: 'JPEREZ', contrasena: 'MiClave123' }",
    "",
    "  2. Rate Limiting: Si hay 10+ intentos fallidos en 10 min → IP bloqueada 5 min",
    "",
    "  3. Joi valida: usuario mínimo 2 caracteres, contraseña mínimo 6",
    "",
    "  4. Busca en SQLite: SELECT * FROM usuarios WHERE usuario = 'JPEREZ'",
    "",
    "  5. Compara contraseña con bcrypt (12 rondas de sal, hash irreversible)",
    "",
    "  6. Si es correcto → Crea sesión en SQLite + cookie 'panelamcham.sid'",
    "     + genera JWT access token (1 hora) + refresh token (7 días)",
    "",
    "  7. Carga TODOS los permisos del usuario (empresas + módulos + acciones)",
    "",
    "  RESULTADO: Usuario entra con sesión activa y permisos cargados.",
    "  ARCHIVOS: src/routes/auth.js | src/middleware/auth.js",
])

content_slide("Cookies, Sesiones y Tokens JWT", [
    "COOKIE (pulsera de acceso):",
    "  Nombre:   'panelamcham.sid'",
    "  httpOnly: true  → JavaScript del navegador NO puede leerla (anti-XSS)",
    "  secure:   true  → Solo funciona con HTTPS en producción",
    "  maxAge:   30 min → Si no haces nada en 30 min, se cierra la sesión",
    "  rolling:  true  → Se renueva con cada click (no expira mientras trabajas)",
    "  store:    SQLite → Las sesiones se guardan en BD, NO en memoria del servidor",
    "",
    "TOKEN JWT:",
    "  Access Token:  duración 1 hora, va en header Authorization: Bearer eyJhbG...",
    "  Refresh Token: duración 7 días, permite renovar el access sin re-login",
    "  POST /api/auth/refresh → envía refreshToken, recibe nuevo par de tokens",
    "",
    "CÓMO SE MANTIENE LA SESIÓN:",
    "  Login 09:00 → cookie expira 09:30",
    "  Click 09:15 → se extiende a 09:45 (rolling)",
    "  Sin actividad 30 min → sesión cerrada, debe re-loguearse",
    "",
    "FRONTEND (sesion.js): window.Sesion guarda tokens en localStorage",
    "  Sesion.headersAutenticacion() → { Authorization: 'Bearer ...', 'X-Empresa-Activa': id }",
])

# ============================================================
# PARTE 5: PERMISOS
# ============================================================
section_slide(5, "Sistema de Permisos (3 niveles)")

content_slide("Permisos: Nivel 1 (Usuario) y Nivel 2 (Empresa)", [
    "NIVEL 1 — PERMISOS GENERALES DEL USUARIO (tabla: usuarios)",
    "  • es_admin_global: acceso total a todo el sistema",
    "  • puede_agregar: puede crear nuevos registros",
    "  • puede_modificar: puede editar registros existentes",
    "  • puede_eliminar: puede borrar registros",
    "",
    "NIVEL 2 — ACCESO POR EMPRESA (tabla: permisos_modulo)",
    "  Un usuario puede tener acceso a una o varias empresas:",
    "  • EMPRESA01 → CIUDAD DE MÉXICO (CDMX)",
    "  • EMPRESA02 → GUADALAJARA (GDL)",
    "  • EMPRESA03 → NORESTE (MTY)",
    "  • EMPRESA04 → NOROESTE (NW)",
    "  • EMPRESA09-12 → Alias comparativos (mapean a EMPRESA01-04)",
    "",
    "  Cada usuario tiene permisos INDEPENDIENTES por empresa.",
    "  Ejemplo: Juan puede ver CDMX y GDL, pero no MTY ni NW.",
    "",
    "ADMIN GLOBAL: Salta todas las verificaciones, puede hacer todo en todo.",
])

content_slide("Permisos: Nivel 3 (Módulo + Acción)", [
    "NIVEL 3 — PERMISOS POR MÓDULO (tabla: permisos_modulo)",
    "  Cada usuario tiene permisos específicos por empresa Y módulo:",
    "",
    "  Ejemplo: Juan Pérez en EMPRESA01 (CDMX):",
    "  ┌────────────────┬─────┬────────────┬─────────┬──────────┐",
    "  │ Módulo         │ Ver │ Cargar/Gda │ Revisar │ Aprobar  │",
    "  ├────────────────┼─────┼────────────┼─────────┼──────────┤",
    "  │ RH             │ SÍ  │ SÍ         │ NO      │ NO       │",
    "  │ Finanzas       │ SÍ  │ NO         │ SÍ      │ NO       │",
    "  │ Eventos        │ SÍ  │ SÍ         │ SÍ      │ NO       │",
    "  │ RESUMEN        │ SÍ  │ NO         │ NO      │ NO       │",
    "  │ SUMMARY        │ NO  │ NO         │ NO      │ NO       │",
    "  └────────────────┴─────┴────────────┴─────────┴──────────┘",
    "",
    "NIVEL ADICIONAL — PERMISOS POR CAPÍTULO (tabla: permisos_edicion_capitulo)",
    "  Controla quién puede EDITAR plantillas por capítulo:",
    "  • CIUDAD DE MÉXICO, GUADALAJARA, NORESTE, NOROESTE, DEFAULT",
    "",
    "VERIFICACIÓN: Cada petición al servidor valida con requireAuth +",
    "  tienePermisoModulo(mapa, empresaId, modulo, accion)",
])

# ============================================================
# PARTE 6: APIS
# ============================================================
section_slide(6, "¿Qué es una API? + Ejemplo real")

content_slide("¿Qué es una API?", [
    "API = el MENSAJERO entre lo que ves en pantalla y donde están los datos.",
    "",
    "  TÚ (navegador)  →  API (mensajero)  →  SERVIDOR (procesa)  →  BASE DE DATOS",
    "",
    "Cómo funciona en SummaCham:",
    "  1. Haces click en 'RESUMEN' en la app",
    "  2. El frontend automáticamente envía esta petición HTTP:",
    "     GET http://localhost:3005/api/reportes/resumen?empresaId=EMPRESA01&anio=2025&mes=6",
    "  3. La petición viaja con tu cookie de sesión y token JWT",
    "  4. El servidor la recibe, verifica permiso, y busca los datos",
    "  5. Responde con un JSON (datos estructurados)",
    "  6. El frontend pinta la tabla con esos datos",
    "",
    "CADA PETICIÓN LLEVA:",
    "  Cookie:  panelamcham.sid=abc123...         (identifica tu sesión)",
    "  Header:  Authorization: Bearer eyJhbG...   (tu token JWT)",
    "  Header:  X-Empresa-Activa: EMPRESA01       (empresa seleccionada)",
    "",
    "MÉTODOS HTTP:  GET (leer) | POST (crear) | PUT (actualizar) | DELETE (eliminar) | PATCH (parcial)",
])

content_slide("Ejemplo real: Abrir el RESUMEN paso a paso", [
    "  1.  Abres RESUMEN.html en la app",
    "  2.  resumen-view.js lee contexto persistido (localStorage: anio, mes, empresa)",
    "  3.  Envía: GET /api/reportes/resumen?empresaId=EMPRESA01&anio=2025&mes=6",
    "  4.  src/server.js dirige /api/reportes/* → src/routes/reportes.js",
    "  5.  middleware/auth.js verifica cookie + JWT → carga mapaPermisos",
    "  6.  reportes.js valida parámetros con Joi",
    "  7.  Llama a planeacionReportesEngine.generarReporte('RESUMEN', ...)",
    "  8.  Engine carga layout de SQLite: qué cuentas y operaciones tiene RESUMEN",
    "  9.  Engine consulta Firebird: saldos de SALDOS25 y presupuestos de PRESUP25",
    "  10. Engine calcula: Real, Presupuesto, Acumulados, Variaciones",
    "  11. Engine ordena por orden_presentacion (flat list, sin agrupar)",
    "  12. Servidor responde JSON → resumen-view.js pinta la tabla",
    "",
    "ARCHIVOS INVOLUCRADOS:",
    "  Frontend: vistas/js/resumen-view.js (~6,600 líneas)",
    "  Ruta:     src/routes/reportes.js",
    "  Motor:    src/services/reportes/planeacionReportesEngine.js",
    "  Firebird: src/services/planeacionCuentasService.js",
    "  Layout:   src/services/layoutService.js",
])

# ============================================================
# PARTE 7: AUTH APIs
# ============================================================
section_slide(7, "APIs de Autenticación")

content_slide("APIs: /api/auth (Login, Refresh, Logout)", [
    "POST /api/auth/login",
    "  Qué hace:  Iniciar sesión",
    "  Envía:     { 'usuario': 'JPEREZ', 'contrasena': 'MiClave123' }",
    "  Responde:  { token, refreshToken,",
    "               usuario: { id, nombre, esAdmin },",
    "               permisos: { EMPRESA01: { RH: {ver:true}, ... } } }",
    "  Errores:   400 (datos inválidos) | 401 (credenciales incorrectas) | 429 (bloqueado)",
    "",
    "POST /api/auth/refresh",
    "  Qué hace:  Renovar token sin re-login (cuando el access token expira)",
    "  Envía:     { 'refreshToken': 'eyJhbG...' }",
    "  Responde:  { token: 'nuevoToken...', refreshToken: 'nuevoRefresh...' }",
    "",
    "POST /api/auth/logout",
    "  Qué hace:  Cerrar sesión (destruye cookie y sesión del servidor)",
    "  Responde:  { mensaje: 'Sesión cerrada' }",
    "",
    "PROTECCIÓN: Rate limiting → 10 intentos / 10 min, bloqueo 5 min",
    "ARCHIVO: src/routes/auth.js",
])

# ============================================================
# PARTE 8: REPORTES
# ============================================================
section_slide(8, "APIs de Reportes")

two_col_slide("APIs: /api/reportes", [
    "GET /api/reportes/summary",
    "  Genera reporte SUMMARY (inglés)",
    "  Params: empresaId, anio, mes",
    "  Responde: { layout: [...],",
    "    cuentas, operaciones, totales }",
    "",
    "GET /api/reportes/resumen",
    "  Genera reporte RESUMEN (español)",
    "  Misma estructura que summary",
    "  pero con layout de RESUMEN",
    "",
    "GET /api/reportes/diagnostico-firebird",
    "  Prueba conexión a Firebird",
    "  Responde: { conectado: bool,",
    "    mensaje, tiempoMs }",
], [
    "POST /api/reportes/operativo-excel-native",
    "  Genera Excel con datos operativos",
    "  Envía: archivo Excel (.xlsx) como",
    "    buffer binario (hasta 20MB)",
    "  Params: nombreArchivo, empresa,",
    "    mes, anio, dataSheetName,",
    "    chartsSheetName, chartMode",
    "  Responde: archivo Excel modificado",
    "",
    "POST /api/reportes/resumen-excel-native",
    "  Genera Excel ejecutivo del resumen",
    "  Misma estructura que operativo",
    "",
    "PARÁMETROS COMUNES:",
    "  empresaId = EMPRESA01..04",
    "  anio = 2024, 2025, 2026",
    "  mes = 1-12",
])

# ============================================================
# PARTE 9: PRESUPUESTOS
# ============================================================
section_slide(9, "APIs de Presupuestos")

content_slide("APIs: /api/presupuestos (lectura y estado)", [
    "GET /api/presupuestos?empresaId=EMPRESA01&anio=2025",
    "  Trae cuentas de presupuesto con sus 12 meses desde Firebird",
    "  Responde: { empresa, anio, periodos, cuentas: [{ numCta, ene, feb, ..., dic }] }",
    "",
    "GET /api/presupuestos/anios?empresaId=EMPRESA01",
    "  Lista años disponibles en Firebird: { anios: [2023, 2024, 2025, 2026] }",
    "",
    "GET /api/presupuestos/estado?modulo=RH&anio=2025&empresaId=EMPRESA01",
    "  Estado actual del presupuesto:",
    "  { estado: 'REVISADO', actualizadoEn: '2025-06-15', actualizadoPor: 'JPEREZ',",
    "    historial: [{ estado, fecha, usuario }, ...] }",
    "  ESTADOS: SIN_CARGAR | EDITANDO | REVISADO | APROBADO | GUARDADO | RECHAZADO",
    "",
    "POST /api/presupuestos/estado  { modulo, anio, accion }",
    "  Transiciones válidas (máquina de estados):",
    "  • 'cargar':    SIN_CARGAR/GUARDADO → EDITANDO  (requiere 'Cargar y guardar')",
    "  • 'revisar':   EDITANDO → REVISADO             (requiere 'Revisar')",
    "  • 'autorizar': REVISADO → APROBADO              (requiere 'Aprobar')",
    "  • 'guardar':   APROBADO → GUARDADO              (requiere 'Aprobar')",
    "  Error 409 si la transición no es válida para el estado actual.",
])

content_slide("APIs: /api/presupuestos (guardar y consolidar)", [
    "POST /api/presupuestos/guardar",
    "  Registra presupuesto guardado oficialmente",
    "  Envía: { empresaId, modulo, anio, datos: [...] }",
    "  Guarda en SQLite (tabla presupuestos_guardados) como respaldo",
    "",
    "POST /api/presupuestos/actualizar-consolidados",
    "  Actualiza cuentas de presupuesto consolidadas en Firebird (PRESUP[yy])",
    "  Envía: { anio, cuentas: [{ numCta, valores: { ene, feb, ..., dic } }] }",
    "  Comportamiento especial:",
    "    • Normaliza numCta a formato canónico de 21 caracteres",
    "    • Actualiza cada mes individualmente en PRESUP[yy]",
    "    • Recalcula cuentas PADRE (450-000-000-00 ingresos, 950-000-000-00 gastos)",
    "    • Consolida desde cuentas detalle (450-001/002/003, 950-001/002/003)",
    "",
    "GET /api/presupuestos/totales-capitulo?empresaId=EMPRESA01&anio=2025",
    "  Totales directos de la tabla PRESUP en Firebird",
    "  Responde: { empresaId, empresa, anio, totales: { income, expense } }",
    "",
    "ARCHIVO: src/routes/presupuestos.js (24KB)",
])

# ============================================================
# PARTE 10: BORRADORES
# ============================================================
section_slide(10, "APIs de Borradores (Draft System)")

content_slide("Flujo completo de Borradores", [
    "Los borradores manejan el ciclo de vida del presupuesto antes de guardarse en COI:",
    "",
    "  ┌──────────┐    ┌──────────┐    ┌──────────┐    ┌──────────┐    ┌──────────┐",
    "  │ EDITANDO │ →  │ PENDIENTE│ →  │ APROBADO │ →  │ GUARDADO │    │RECHAZADO │",
    "  │ (editar) │    │ (enviar) │    │(autorizar│    │(a COI)   │    │(regresa  │",
    "  │          │    │          │    │          │    │          │    │ a editar)│",
    "  └──────────┘    └──────────┘    └──────────┘    └──────────┘    └──────────┘",
    "",
    "POST /api/borradores/guardar    → Crea o actualiza borrador (autoguardado)",
    "POST /api/borradores/enviar     → EDITANDO → PENDIENTE (si admin: auto-aprueba)",
    "POST /api/borradores/revisar    → EDITANDO → REVISADO (o cancelar: PENDIENTE → EDITANDO)",
    "POST /api/borradores/autorizar  → PENDIENTE → APROBADO (requiere permiso 'Aprobar')",
    "POST /api/borradores/rechazar   → Cualquier estado → RECHAZADO (con motivo)",
    "POST /api/borradores/finalizar  → APROBADO → GUARDADO (escribe en Firebird COI)",
    "POST /api/borradores/descartar  → Elimina borrador completamente (hard delete)",
    "",
    "GET  /api/borradores/listar     → Lista borradores (filtrado por permisos)",
    "GET  /api/borradores/estado     → Estado actual + progreso recontabilización",
    "GET  /api/borradores/historial  → Historial con filtros (fecha, estado, acción, usuario)",
    "GET  /api/borradores/detalle/:id → Datos completos del borrador",
])

content_slide("Recontabilización y Guardar en COI", [
    "RECONTABILIZACIÓN (recalcular cuentas padre desde hijas):",
    "",
    "POST /api/borradores/recontabilizar",
    "  Envía: { empresaId: 'EMPRESA01', anio: 2025 }",
    "  Recalcula TODAS las cuentas padre de la empresa/año en PRESUP[yy]",
    "  Ejemplo: cuenta padre 450-000-000-00 = suma de 450-001 + 450-002 + 450-003",
    "  Registra progreso en tabla recontabilizacion_progreso",
    "",
    "POST /api/borradores/recontabilizar/todas  (SOLO ADMIN)",
    "  Recontabiliza TODAS las empresas para un año",
    "  Responde: { exitosas: 4, fallidas: 0, resultados: [...] }",
    "",
    "FINALIZAR (Guardar en COI):",
    "POST /api/borradores/finalizar  { borradorId: 123 }",
    "  1. Verifica que el borrador esté en estado APROBADO",
    "  2. Verifica permiso 'Cargar y guardar' del usuario",
    "  3. Llama a guardarAutorizado() → escribe en tabla PRESUP[yy] de Firebird",
    "  4. Cambia estado a GUARDADO, registra ejecutor y fecha",
    "  5. Envía notificación al autor original",
    "",
    "ARCHIVO: src/routes/borradores.js (32KB, el más grande del sistema)",
])

# ============================================================
# PARTE 11: USUARIOS
# ============================================================
section_slide(11, "APIs de Usuarios y Perfil")

two_col_slide("APIs: /api/usuarios", [
    "GET /api/usuarios  (Solo Admin)",
    "  Lista todos los usuarios",
    "  Params: limite (default 100), offset",
    "",
    "GET /api/usuarios/:id",
    "  Datos de un usuario específico",
    "  (Admin o el propio usuario)",
    "",
    "POST /api/usuarios/crear (Solo Admin)",
    "  Body validado con Joi:",
    "  { usuario (2-32 chars),",
    "    nombres (1-80 chars),",
    "    apellidoPrimero (1-120 chars),",
    "    correo (email válido),",
    "    contrasena (min 8 chars),",
    "    esAdminGlobal (boolean),",
    "    permisosGenerales: { ... },",
    "    permisos: { EMPRESA01: { RH: {...} } }",
    "  }",
], [
    "PUT /api/usuarios/:id",
    "  Editar usuario (Admin o self)",
    "  Contraseña es opcional",
    "",
    "DELETE /api/usuarios/:id (Admin)",
    "  Elimina usuario + permisos (CASCADE)",
    "",
    "PATCH /api/usuarios/:id/restablecer-contrasena",
    "  Reset contraseña (Solo Admin)",
    "",
    "GET /api/usuarios/:id/permisos",
    "  Ver permisos de un usuario",
    "",
    "PUT /api/usuarios/:id/permisos",
    "  Actualizar matriz de permisos",
    "  (Borra los viejos e inserta nuevos)",
    "",
    "--- PERFIL (/api/perfil) ---",
    "GET  /api/perfil → Mi perfil completo",
    "PATCH /api/perfil/password",
    "  { contrasenaActual, nuevaContrasena }",
])

# ============================================================
# PARTE 12: LAYOUTS
# ============================================================
section_slide(12, "APIs de Plantillas (Layouts)")

content_slide("APIs: /api/layouts-config (cargar y guardar)", [
    "BASE: /api/layouts-config    ARCHIVO: src/routes/layoutRoutes.js",
    "",
    "CARGAR:",
    "  GET  /.../RESUMEN/2025/DEFAULT          → Layout completo (cuentas + operaciones)",
    "  GET  /.../RESUMEN/anios                  → Qué años existen",
    "  GET  /.../RESUMEN/2025/capitulos         → Qué capítulos hay (CDMX, GDL, etc.)",
    "  GET  /.../RESUMEN/2025/estadisticas      → Cuántas cuentas y operaciones tiene",
    "  GET  /.../RESUMEN/2025/existe            → Verificar si ya existe este layout",
    "",
    "GUARDAR (2 llamadas separadas para cuentas y operaciones):",
    "  POST /.../RESUMEN/2025/DEFAULT/cuentas   → Guardar todas las cuentas",
    "      Body: { empresaId, cuentas: [{ CUENTA, NOMBRE, CAPITULO,",
    "              'SECCION Principal', 'SECCION Secundaria',",
    "              operacion_factor, orden_presentacion, visible }] }",
    "",
    "  POST /.../RESUMEN/2025/operaciones       → Guardar todas las operaciones",
    "      Body: { empresaId, capitulo, operaciones: [{ Clase, SECCION,",
    "              'sum-row', 'result-row', 'net-row', formula_terms,",
    "              signos, orden_presentacion, visible }] }",
    "",
    "NOTA: El guardado es TRANSACCIONAL → todo o nada (SQLite transaction)",
])

content_slide("APIs: /api/layouts-config (edición y administración)", [
    "EDICIÓN INDIVIDUAL:",
    "  PUT    /.../RESUMEN/2025/DEFAULT/cuenta/4100    → Editar cuenta 4100",
    "  DELETE /.../RESUMEN/2025/DEFAULT/cuenta/4100    → Eliminar cuenta 4100",
    "  PUT    /.../RESUMEN/2025/operacion/TOTAL_INCOME → Editar una operación",
    "  DELETE /.../RESUMEN/2025/operacion/TOTAL_INCOME → Eliminar una operación",
    "",
    "ADMINISTRACIÓN:",
    "  POST   /.../RESUMEN/copiar                      → Copiar layout de otro año",
    "      Body: { origen: 2024, destino: 2025 }  (transaccional, borra destino primero)",
    "  POST   /.../RESUMEN/2025/DEFAULT/reordenar      → Cambiar orden de filas",
    "      Body: [{ cuenta: '4100', orden_presentacion: 10 }, ...]",
    "  DELETE /.../RESUMEN/2025                         → Eliminar layout completo del año",
    "  DELETE /.../RESUMEN/2025/DEFAULT                 → Eliminar solo un capítulo",
    "  GET    /.../RESUMEN/2025/DEFAULT/bitacora        → Ver registro de cambios",
    "",
    "SECCIONES:",
    "  POST   /.../seccion/crear       → Crear nueva sección",
    "  PUT    /.../seccion/renombrar   → Renombrar (actualiza cuentas y secciones en cascada)",
    "",
    "INSERCIÓN (/api/insercion):",
    "  POST /api/insercion/agregar-fila    → Agregar fila con validación anti-duplicados",
    "  POST /api/insercion/agregar-seccion → Agregar sección con validación jerárquica",
])

# ============================================================
# PARTE 13: SALDOS Y OTROS
# ============================================================
section_slide(13, "APIs: Saldos, Módulos, Notificaciones y más")

two_col_slide("APIs: Saldos, Módulos, Config", [
    "SALDOS (/api/saldos):",
    "  GET /api/saldos",
    "    Saldos de cuentas desde Firebird",
    "    Params: anio, empresaId, cuentas",
    "    Resp: 12 meses de cargo/abono",
    "",
    "  GET /api/saldos/anios",
    "    Años disponibles en Firebird",
    "",
    "  GET /api/saldos/cuentas",
    "    Cuentas por año (filtrable)",
    "",
    "  GET /api/saldos/catalogo",
    "    Catálogo completo de cuentas",
    "",
    "MÓDULOS (/api/modulos):",
    "  GET /api/modulos → Listar todos",
    "  GET /api/modulos/:id → Detalle",
    "  POST /api/modulos/summary-resumen-e",
    "    Detalle con comparación presupuesto",
    "  GET /api/modulos/summary-anios",
    "    Años desde tablas SALDOS",
], [
    "NOTIFICACIONES (/api/notificaciones):",
    "  GET  /  → Lista alertas (límite 20)",
    "  PATCH /limpiar → Marcar todas leídas",
    "  PATCH /:id/leida → Marcar una leída",
    "",
    "COMENTARIOS (/api/comentarios):",
    "  CRUD para comentarios en celdas",
    "  de presupuesto (por celda_id)",
    "",
    "GRÁFICAS (/api/graficas-config):",
    "  GET  / → Cargar config guardada",
    "  POST / → Guardar config (admin)",
    "",
    "FIREBIRD CONFIG:",
    "  GET  / → Ver config actual",
    "  PUT  /base → Cambiar host/puerto",
    "  PUT  /empresa/:id → Cambiar ruta BD",
    "  POST /probar/:id → Test conexión",
    "",
    "EMPRESAS (/api/empresas):",
    "  GET / → Lista de empresas disponibles",
])

# ============================================================
# PARTE 14: MAPA APIS
# ============================================================
section_slide(14, "Mapa completo de TODAS las APIs")

two_col_slide("80+ Endpoints organizados por grupo", [
    "AUTENTICACIÓN (3):",
    "  POST /api/auth/login",
    "  POST /api/auth/refresh",
    "  POST /api/auth/logout",
    "",
    "USUARIOS (8):",
    "  GET|POST /api/usuarios",
    "  GET|PUT|DEL /api/usuarios/:id",
    "  PUT /api/usuarios/:id/permisos",
    "  PATCH /:id/restablecer-contrasena",
    "",
    "PERFIL (2):",
    "  GET /api/perfil",
    "  PATCH /api/perfil/password",
    "",
    "REPORTES (5):",
    "  GET /api/reportes/summary",
    "  GET /api/reportes/resumen",
    "  GET /api/reportes/diagnostico-firebird",
    "  POST /api/reportes/operativo-excel-native",
    "  POST /api/reportes/resumen-excel-native",
], [
    "PRESUPUESTOS (6):",
    "  GET estado | anios | totales-capitulo",
    "  POST estado | guardar | actualizar",
    "",
    "BORRADORES (12):",
    "  GET listar|estado|historial|detalle|progreso",
    "  POST guardar|enviar|revisar|autorizar",
    "  POST rechazar|finalizar|descartar",
    "  POST recontabilizar|recontabilizar/todas",
    "",
    "LAYOUTS (15+):",
    "  GET|POST|PUT|DEL layouts-config/...",
    "  POST insercion/agregar-fila|seccion",
    "",
    "SALDOS (4): GET saldos|anios|cuentas|catalogo",
    "MÓDULOS (4): GET|POST modulos/...",
    "GRÁFICAS (2): GET|POST graficas-config",
    "BACKUPS (6): GET|POST|PUT backups/...",
    "NOTIFICACIONES (3): GET|PATCH",
    "FIREBIRD CONFIG (4): GET|PUT|POST",
    "EMPRESAS (1): GET",
    "COMENTARIOS (4+): CRUD",
])

# ============================================================
# PARTE 15: FIREBIRD
# ============================================================
section_slide(15, "Firebird: De dónde salen los números")

content_slide("Conexión a Firebird (Aspel COI)", [
    "Firebird es la base de datos del sistema contable Aspel COI.",
    "SummaCham se CONECTA a Firebird para LEER los datos reales.",
    "",
    "CONFIGURACIÓN DE CONEXIÓN (src/services/firebirdService.js):",
    "  • host: IP del servidor (localhost o remoto)",
    "  • port: 3050 (default Firebird)",
    "  • database: ruta al archivo .fdb de cada empresa",
    "  • user/password: credenciales del servidor Firebird",
    "",
    "DETECCIÓN AUTOMÁTICA LOCAL vs REMOTA:",
    "  LOCAL  → host = localhost/127.0.0.1, timeout 3s, 0 reintentos, query timeout 10s",
    "  REMOTA → cualquier otro host, timeout 120s, 3 reintentos, query timeout 120s",
    "",
    "TABLAS POR AÑO (Firebird tiene tablas separadas por año fiscal):",
    "  2024: CUENTAS24, SALDOS24, PRESUP24",
    "  2025: CUENTAS25, SALDOS25, PRESUP25",
    "  2026: CUENTAS26, SALDOS26, PRESUP26",
    "",
    "PERFORMANCE: Queries lentas (>2s) se loguean automáticamente con duración y filas",
    "ARCHIVO: src/services/firebirdService.js → ejecutarConsulta(empresaId, sql, params)",
])

content_slide("Estructura de las tablas Firebird", [
    "CUENTAS[yy] — Catálogo de cuentas contables:",
    "  NUM_CTA:     Número de cuenta (ej: '401-000-000-00')",
    "  NOMBRE:      Descripción (ej: 'Membership')",
    "  NIVEL:       1=mayor, 2=sub, 3=detalle, 4=auxiliar",
    "  NATURALEZA:  'D'=Deudora(gasto) o 'A'=Acreedora(ingreso)",
    "  STATUS:      'A'=Activa (solo se consultan activas)",
    "",
    "SALDOS[yy] — Movimientos mensuales:",
    "  NUM_CTA:     Cuenta",
    "  EJERCICIO:   Año fiscal",
    "  CARGO01..12: Dinero que ENTRA (débitos) por mes",
    "  ABONO01..12: Dinero que SALE (créditos) por mes",
    "",
    "PRESUP[yy] — Presupuestos por mes:",
    "  NUM_CTA:       Cuenta",
    "  EJERCICIO:     Año fiscal",
    "  PRESUP01..12:  Monto presupuestado por mes",
    "",
    "RESOLUCIÓN DE EMPRESA → BASE DE DATOS:",
    "  EMPRESA01 → ruta/empresa1.fdb",
    "  EMPRESA02 → ruta/empresa2.fdb",
    "  EMPRESA09 → alias a EMPRESA01 (comparativos)",
])

# ============================================================
# PARTE 16: CALCULOS
# ============================================================
section_slide(16, "Cálculos: Real, Presupuesto, Acumulado")

content_slide("Cómo se calcula el REAL (dato contable actual)", [
    "La NATURALEZA de la cuenta decide la fórmula:",
    "",
    "CUENTA ACREEDORA (ingresos) — cuentas que empiezan con 2, 3, 4:",
    "  Real = ABONOS - CARGOS",
    "  Ejemplo: Cuenta 4100 'Membership' en Junio:",
    "    ABONO06 = 150,000 (créditos)   CARGO06 = 15,000 (débitos)",
    "    REAL = 150,000 - 15,000 = 135,000",
    "",
    "CUENTA DEUDORA (gastos) — cuentas que empiezan con 1, 5-9:",
    "  Real = CARGOS - ABONOS",
    "  Ejemplo: Cuenta 5100 'Salaries' en Junio:",
    "    CARGO06 = 120,000 (débitos)   ABONO06 = 5,000 (créditos)",
    "    REAL = 120,000 - 5,000 = 115,000",
    "",
    "¿POR QUÉ ES DIFERENTE?",
    "  En contabilidad, un ingreso AUMENTA con abonos (créditos)",
    "  y un gasto AUMENTA con cargos (débitos). La fórmula ajusta",
    "  para que el número siempre sea POSITIVO cuando la cuenta",
    "  tiene saldo a favor.",
    "",
    "ARCHIVO: src/services/planeacionCuentasService.js → determinarNaturalezaReal()",
])

content_slide("Presupuesto y Acumulados (YTD)", [
    "PRESUPUESTO (se lee directo, sin cálculo):",
    "  SELECT PRESUP06 FROM PRESUP25 WHERE NUM_CTA = '4100'",
    "  Resultado: 140,000 → se usa TAL CUAL, no necesita fórmula",
    "",
    "ACUMULADOS (YTD = Year-To-Date):",
    "  Si seleccionas MES = 6 (Junio):",
    "  Real Acumulado  = Real Ene + Real Feb + ... + Real Jun",
    "  Ppto Acumulado  = PRESUP01 + PRESUP02 + ... + PRESUP06",
    "  Real Acum AA    = Misma suma pero del año ANTERIOR (SALDOS24)",
    "",
    "LOS 6 VALORES QUE VES EN CADA FILA DEL REPORTE:",
    "  ┌───────────────────┬───────────────────┬───────────────────┐",
    "  │ actualMonth (Real)│ planMonth (Ppto)  │ prevMonth (AA)    │",
    "  │ actualYTD (Acum)  │ planYTD (Acum)    │ prevYTD (Acum AA) │",
    "  └───────────────────┴───────────────────┴───────────────────┘",
    "",
    "VARIACIONES (calculadas en el frontend):",
    "  Var vs Plan = (Real - Ppto) / Ppto × 100  →  Ejemplo: -3.6%",
    "  Var vs AA   = (Real - Real AA) / Real AA × 100",
    "",
    "ARCHIVOS: planeacionCuentasService.js (SQL) + planeacionReportesEngine.js (cálculo)",
])

content_slide("Normalización de cuentas (formato canónico)", [
    "Las cuentas se convierten a un formato canónico de 21 caracteres para comparación:",
    "",
    "normalizarCuentaCanonica(valor):",
    "  1. Extrae solo dígitos del número de cuenta",
    "  2. Rellena con ceros hasta 11 caracteres",
    "  3. Agrega sufijo de nivel (1-4)",
    "  Resultado: 21 caracteres fijos para comparación exacta",
    "",
    "EJEMPLOS:",
    "  '401'           → '40100000000' + '0000000000' + '1' = 21 chars",
    "  '401-000-000-00' → mismo resultado (los guiones se eliminan)",
    "  '5100'          → '51000000000' + '0000000000' + '1'",
    "",
    "cuentaVisibleDesdeCanonica(canonica):",
    "  Convierte de vuelta al formato legible: '40100000000' → '401-000-000-00'",
    "",
    "¿POR QUÉ ES NECESARIO?",
    "  Firebird puede tener '4100' y SQLite '401-000-000-00'",
    "  para la MISMA cuenta. La normalización permite que se",
    "  emparejen correctamente al cruzar datos de ambas bases.",
    "",
    "ARCHIVO: src/services/reportes/planeacionReportesEngine.js",
])

# ============================================================
# PARTE 17: MOTOR
# ============================================================
section_slide(17, "Motor de Reportes (Engine)")

content_slide("construirReporteResumen() — el corazón del sistema", [
    "Esta función en planeacionReportesEngine.js construye TODO el reporte.",
    "",
    "ENTRADA:",
    "  • definiciones: array de cuentas del layout (de SQLite)",
    "  • configAgrupacion: operaciones/fórmulas del layout",
    "  • capituloSeleccionado: filtro de capítulo",
    "  • planeacionData: datos reales de Firebird (6 valores por cuenta)",
    "  • opciones: { permitirFormulaSecciones, ordenManualTotal, includeEmptyConfiguredSections }",
    "",
    "PROCESO (paso a paso):",
    "  1. Parsear configuración → extraer sum-row, result-row, net-row de operaciones",
    "  2. Construir mapa de principales → agrupar cuentas por 'SECCION Principal'",
    "  3. Construir mapa de secciones → agrupar dentro de cada principal por 'SECCION Secundaria'",
    "  4. Aplicar fórmulas manuales → si permitirFormulaSecciones, calcular totales custom",
    "  5. Construir filas de agregación → consolidados, operativos, resultado, neto, final",
    "  6. Intercalar operaciones libres → colocar operaciones cerca de sus secciones",
    "  7. Aplicar orden global → si ordenManualTotal, ordenar TODO por orden_presentacion",
    "",
    "SALIDA: { principals: [...], layout: [...flat array...] }",
])

content_slide("Estructura del layout generado (tipos de fila)", [
    "El layout es un ARRAY PLANO de filas tipadas:",
    "",
    "  type:'principal'  → Encabezado de sección principal (ej: 'Income (CDMX)')",
    "    Campos: label, order, totals {actualMonth, planMonth, prevMonth, ...YTD}",
    "",
    "  type:'secundaria' → Sub-sección (ej: 'Membership Income')",
    "    Campos: label, order, totals, cuentas: [...]",
    "",
    "  type:'cuenta'     → Fila individual de cuenta contable",
    "    Campos: label, nombre, cuenta, order, totals, factor",
    "",
    "  type:'operation'  → Operación con fórmula manual",
    "    Campos: label, order, manualFormula, totals",
    "",
    "  type:'group'      → Suma/consolidado de secciones (ej: 'CONSOLIDATED INCOME')",
    "    Campos: label, order, principals: [...], operaciones: [...]",
    "",
    "  type:'result'     → Resultado (ej: 'OPERATING RESULTS')",
    "  type:'net'        → Neto (ej: 'NET INCOME')",
    "  type:'final'      → Resultado final (ej: 'RESULTADO NETO CONSOLIDADO')",
    "",
    "ORDEN: Todo el array se ordena por orden_presentacion (flat, sin agrupación).",
    "  Las cuentas y operaciones se mezclan libremente según su orden.",
])

content_slide("Cómo funcionan las Operaciones (fórmulas)", [
    "Las operaciones son filas que CALCULAN totales usando fórmulas guardadas en JSON:",
    "",
    "EJEMPLO 1 — Suma de secciones:",
    "  Operación: 'TOTAL INCOME'",
    "  formula_terms: [",
    "    { operator: '+', type: 'section', value: 'Income (CDMX)' },",
    "    { operator: '+', type: 'section', value: 'Income (GDL)' },",
    "    { operator: '+', type: 'section', value: 'Income (MTY)' }",
    "  ]",
    "  Resultado: TOTAL INCOME = suma de TODOS los ingresos de 3 ciudades",
    "",
    "EJEMPLO 2 — Referencia recursiva a otras operaciones:",
    "  Operación: 'NET INCOME'",
    "  formula_terms: [",
    "    { operator: '+', type: 'operation', value: 'CONSOLIDATED INCOME' },",
    "    { operator: '-', type: 'operation', value: 'CONSOLIDATED EXPENSE' }",
    "  ]",
    "",
    "TIPOS DE TÉRMINOS: section | operation | account | constant",
    "OPERADORES: + | - | * | /",
    "FUNCIÓN: calcularTotalesOperacion() evalúa recursivamente cada término",
])

# ============================================================
# PARTE 18: GESTOR PLANTILLAS
# ============================================================
section_slide(18, "Gestor de Plantillas (14,700+ líneas)")

content_slide("Gestor de Plantillas — Arquitectura", [
    "El Gestor de Plantillas (plantillas.js) es el archivo MÁS GRANDE del sistema.",
    "Define la ESTRUCTURA de los reportes: sin plantilla, no hay reporte.",
    "",
    "ARCHIVO: vistas/js/plantillas.js → 14,753 líneas de código",
    "",
    "CONSTANTES CLAVE:",
    "  MANUAL_ORDER_ONLY = true     → Siempre orden manual (no auto-inferencia)",
    "  FORCE_EDIT_MODE = true       → Modo edición forzado (desarrollo local)",
    "  AUTO_OPERACIONES_DISABLED = true → No generar operaciones automáticamente",
    "",
    "ESTADO INTERNO (state):",
    "  { modulo: 'RESUMEN',          // Módulo seleccionado",
    "    anio: null,                   // Año seleccionado",
    "    capitulo: null,               // Capítulo (CDMX, GDL, etc.)",
    "    empresaId: null,              // Empresa activa",
    "    layout: null,                 // Layout completo cargado",
    "    cuentas: [],                  // Array de cuentas",
    "    operaciones: [],              // Array de operaciones",
    "    unsavedChanges: false,        // ¿Hay cambios sin guardar?",
    "    editMode: false,              // ¿En modo edición?",
    "    inlineOrderMode: false,       // ¿En modo drag-drop?",
    "    changeLog: [] }              // Registro de cambios",
])

content_slide("Gestor: Flujo de Carga y Renderizado", [
    "FLUJO DE CARGA (al abrir la página):",
    "  1. loadInitialData() → carga años, capítulos",
    "  2. tryLoadLayout() → intenta cargar el layout guardado",
    "  3. loadLayout() → GET /layouts-config/{modulo}/{anio}/{capitulo}",
    "     • Extrae cuentas de layout.cuentas o layout[modulo]",
    "     • Procesa operaciones con extractColumnConfigFromOperations()",
    "     • Hidrata parentSection/parentSubsection con hydrateOperationPlacement()",
    "  4. renderLayout() → llama a buildPreviewRowsForEditor()",
    "",
    "buildPreviewRowsForEditor() — FUNCIÓN CRÍTICA:",
    "  • Construye lista PLANA ordenada por orden_presentacion",
    "  • NO agrupa por secciones (cambio del 2026-02-05)",
    "  • Las cuentas y operaciones se intercalan libremente",
    "  • Detecta transiciones de sección para insertar encabezados",
    "  • Incluye operaciones huérfanas con appendMissingOperations()",
    "  • Retorna array con type: principal | subsection | account | operation",
    "",
    "FLUJO DE GUARDADO:",
    "  saveLayout() → valida → syncVisualOrder() → confirma →",
    "  POST cuentas → POST operaciones → limpia changeLog → addToBitacora()",
])

content_slide("Gestor: Funciones principales", [
    "GESTIÓN DE ELEMENTOS:",
    "  openAddModal()    → Modal para agregar cuentas, operaciones o secciones",
    "  confirmAdd()      → Valida y agrega al array correspondiente",
    "  confirmEdit()     → Edita elemento en panel lateral",
    "  deleteElement()   → Elimina del layout (marca unsavedChanges=true)",
    "",
    "ORDENAMIENTO:",
    "  getTemplateRowsForReorder()    → Extrae filas para el modal de reorden",
    "  applyTemplateRowsOrder(rows)   → Aplica nuevos orden_presentacion",
    "  updateLayoutOrderPanel()       → Panel lateral con contadores",
    "  Drag-drop: drag-and-drop-reorder.js + plantillas-reordenar.js",
    "",
    "MÓDULOS Y EMPRESAS:",
    "  filtrarModulosPorCapitulo()   → Filtra módulos por disponibilidad del capítulo",
    "  obtenerEmpresaIdApi()         → Resuelve empresa activa (state → context → default)",
    "  resolverEmpresaIdContexto()   → Consulta selector global → Sesion → default",
    "",
    "VERIFICACIÓN DE PERMISOS (cada 3 segundos):",
    "  checkAuthState() → Sesion.esAdminGlobal() || verificarPermisoCapitulo()",
    "",
    "CADA LAYOUT ES ÚNICO POR: Empresa + Módulo + Año + Capítulo",
])

# ============================================================
# PARTE 19: GRAFICAS
# ============================================================
section_slide(19, "Sistema de Gráficas")

content_slide("Gráficas: Arquitectura y tipos", [
    "Las gráficas toman los MISMOS datos del reporte RESUMEN y los visualizan.",
    "",
    "TIPOS DE GRÁFICA DISPONIBLES:",
    "  • Barra (bar) — la más usada, comparación entre categorías",
    "  • Línea (line) — tendencias mensuales",
    "  • Pie (pie) — distribución porcentual",
    "  • Dona (doughnut) — similar a pie pero con hueco central",
    "  • Radar (radar) — comparación multidimensional",
    "  • Área (area) — tendencia con relleno",
    "",
    "LOS 3 VALORES QUE SE GRAFICAN (colores por defecto):",
    "  Real Acumulado (actualYTD)    → Azul oscuro #0d47a1",
    "  Ppto Acumulado (planYTD)      → Azul claro  #60a5fa",
    "  Real Acum AA   (prevYTD)      → Gris        #94a3b8",
    "",
    "GRÁFICAS PREDEFINIDAS:",
    "  1. Resultado Operativo por capítulo (CDMX, GDL, MTY, NW)",
    "  2. Resultado Neto por capítulo",
    "  3. Consolidados (Ingreso total vs Gasto total)",
    "  4. Ingreso por capítulo (CDMX/GDL/MTY/NW INCOME)",
    "  5. Ingreso Nacional (Committees, Membership, Events, Services, T&IC)",
    "  6. Gráficas PERSONALIZADAS con Chart Wizard (4 pasos)",
])

content_slide("Gráficas: Sistema de Snapshots y Sidebar", [
    "SISTEMA DE SNAPSHOTS (localStorage):",
    "  Cada vez que se renderiza la tabla RESUMEN, se guarda un 'snapshot'",
    "  en localStorage con la clave: resumen_tabla_snapshot:{empresaId}:{anio}:{mes}",
    "  Contiene: array de { label, totals: { actual, plan, prev, ...YTD, ...var } }",
    "",
    "  Las gráficas LEEN este snapshot, NO hacen otra llamada al API.",
    "  getRowData(snapshotMap, labels) busca por label normalizado (sin acentos, uppercase)",
    "  Si no encuentra exacto, busca variantes (ej: 'CONSOLIDATED' = 'CONSOLIDADO')",
    "",
    "OPERATIVO SIDEBAR (operativo-sidebar.js, 1,633 líneas):",
    "  Panel lateral que muestra gráficas en tiempo real junto a la tabla:",
    "  • Se actualiza con MutationObserver (debounce 120ms)",
    "  • Extrae datos de filas con clase 'sum-row-operativo'",
    "  • Muestra: Ppto Acumulado vs Real Acumulado vs Presupuesto Anual",
    "  • Soporta gráficas custom de GraficasConfig",
    "  • Escucha eventos: tabla-actualizada, presupuesto-editado, graficas-config-updated",
    "",
    "CHART WIZARD (4 pasos para crear gráficas personalizadas):",
    "  1. Seleccionar tipo de gráfica (bar, line, pie, etc.)",
    "  2. Elegir filas/datos del snapshot a incluir",
    "  3. Configurar colores, labels y opciones",
    "  4. Guardar → se persiste en graficas_config (SQLite) via /api/graficas-config",
])

# ============================================================
# PARTE 20: BD, BACKUPS, SEGURIDAD
# ============================================================
section_slide(20, "Base de Datos, Backups y Seguridad")

content_slide("SQLite: Esquema completo de tablas", [
    "Archivo: datos/panel.sqlite — Modo WAL (concurrencia de lectura)",
    "Pragma: journal_mode=WAL, synchronous=NORMAL, cache_size=10000, busy_timeout=5000",
    "",
    "TABLAS PRINCIPALES (15+):",
    "  usuarios              → id, usuario, nombres, apellidos, correo, contrasena (bcrypt)",
    "  permisos_modulo       → usuario_id, empresa_id, modulo, puede_leer/cargar/revisar/aprobar",
    "  permisos_edicion_capitulo → usuario_id, capitulo, puede_editar",
    "  presupuestos_estado   → empresa_id, modulo, anio, estado, actualizado_en/por",
    "  presupuestos_estado_historial → historial de transiciones de estado",
    "  presupuestos_guardados → respaldo de datos guardados oficialmente",
    "  PLAN_BORRADORES       → autoguardado de datos en edición",
    "  layout_cuentas        → cuentas por empresa/modulo/anio/capitulo con orden_presentacion",
    "  layout_operaciones    → fórmulas multi-fila (una fila por operacion_tipo)",
    "  layout_secciones      → agrupaciones (Income, Expense, etc.)",
    "  layout_bitacora       → auditoría de cambios en layouts",
    "  graficas_config       → configuración de gráficas por empresa",
    "  comentarios_celdas    → comentarios en celdas específicas",
    "  notificaciones        → alertas para usuarios",
    "  sesiones              → express-session store en SQLite",
])

content_slide("Migraciones automáticas y Seeding", [
    "MIGRACIONES (al iniciar la app, src/db/sqlite.js):",
    "  El sistema verifica y agrega automáticamente columnas faltantes:",
    "  layout_cuentas:     + visible, orden_presentacion, operacion_factor",
    "  layout_operaciones: + visible, orden_presentacion, formula_json, operacion_etiqueta",
    "  presupuestos_estado: + capitulo",
    "",
    "  Reconstrucción de índices:",
    "  • Si layout_cuentas tiene UNIQUE constraint legacy → rebuild sin constraint",
    "    (permite cuentas duplicadas para almacenamiento multi-fila)",
    "  • Crea índices de lookup: idx_layout_cuentas_lookup(empresa_id, modulo, anio, capitulo)",
    "",
    "SEEDING DE USUARIOS (al primer inicio):",
    "  1. Crea admin ICONET con contraseña de env var o fallback",
    "  2. Lee config/seed_users.json para crear usuarios iniciales",
    "  3. Asigna permisos por empresa/módulo según JSON",
    "  4. Admin históricos (AA, AMB) se promueven automáticamente a admin global",
    "  5. Todos los admin globales reciben permisos completos en TODAS las empresas/módulos",
    "",
    "SEEDING DE LAYOUTS:",
    "  • Desde archivos JSON y Excel (.xlsx) en directorio info-importante/",
    "  • CDMX: layouts de 2005-2025 desde archivos Excel",
    "  • Solo ejecuta si no existen layouts previos (no sobrescribe)",
])

content_slide("Sistema de Backups y Seguridad", [
    "BACKUPS AUTOMÁTICOS (src/services/backupService.js):",
    "  • Al iniciar: primer backup inmediato",
    "  • Cada 60 minutos: backup automático con setInterval",
    "  • Máximo 24 backups (los más antiguos se eliminan)",
    "  • Nombre: panel_2026-02-09_14-30-22.sqlite",
    "  • Verificación: compara tamaño original vs copia + checksum SHA256",
    "  • Restauración: crea backup de seguridad antes de restaurar",
    "",
    "SEGURIDAD (10 capas de protección):",
    "  1.  Contraseñas: bcrypt con 12 rondas de sal (irreversible)",
    "  2.  JWT: Access token 1hr + Refresh 7 días (renovación automática)",
    "  3.  Sesiones: Cookie httpOnly (JS no puede leerla) + SQLite store",
    "  4.  Rate Limiting: 10 intentos/10min, bloqueo 5 min",
    "  5.  Helmet: Cabeceras HTTP seguras (XSS, HSTS, clickjacking)",
    "  6.  CORS: Solo orígenes autorizados (localhost + dominio producción)",
    "  7.  Joi: Validación de TODOS los datos de entrada antes de procesarlos",
    "  8.  SQL seguro: Queries parametrizados (?) previenen SQL injection",
    "  9.  Backups: Copias automáticas cada 60 min con checksum SHA256",
    "  10. Auditoría: Login, cambios, queries lentas, bitácora de layouts",
])

content_slide("Compilación, Distribución y Auto-Update", [
    "COMPILACIÓN:",
    "  npm install                    → Instalar ~120+ dependencias",
    "  npm run dist                   → Crear instalador Windows (NSIS)",
    "  npm run build:portable         → Crear versión portable (.exe)",
    "",
    "ARCHIVOS GENERADOS en ./dist/:",
    "  SummaCham Setup 4.1.0.exe      → Instalador (requiere admin)",
    "  SummaCham 4.1.0.exe            → Portable (sin instalación)",
    "  latest.yml                     → Metadatos para auto-update",
    "",
    "AUTO-ACTUALIZACIÓN (electron-updater):",
    "  1. La app verifica en GitHub Releases si hay versión nueva (al iniciar)",
    "  2. Si hay: pregunta 'Descargar nueva versión?'",
    "  3. Descarga en segundo plano (no interrumpe el trabajo)",
    "  4. Al cerrar la app, instala automáticamente",
    "  5. Los DATOS del usuario se conservan (userData/datos/)",
    "     Solo se actualiza el código de la aplicación",
    "",
    "ESTRUCTURA DE ARCHIVOS PRINCIPALES:",
    "  main.js                → Proceso Electron (ventana + auto-update)",
    "  src/server.js           → Servidor Express (middleware + rutas)",
    "  src/db/sqlite.js        → Inicialización BD + migraciones + seeding",
    "  src/routes/*.js         → 16+ archivos de rutas API",
    "  src/services/*.js       → Lógica de negocio (Firebird, layouts, reportes)",
    "  vistas/                 → HTML + JS + CSS del frontend",
])

# ============================================================
# SLIDE FINAL
# ============================================================
s = title_slide("SummaCham - PanelAMCHAM", "Resumen Técnico Final")
add_multiline(s, 0.6, 3.0, 8.8, 3.0, [
    "Arquitectura en capas: Electron + Express + SQLite + Firebird",
    "80+ endpoints REST documentados en 16+ archivos de rutas",
    "Dual-database: SQLite (config local) + Firebird (datos contables Aspel COI)",
    "6 valores por cuenta: Real, Ppto, Real AA, Real Acum, Ppto Acum, Real AA Acum",
    "Motor de reportes con fórmulas recursivas y orden manual global",
    "Gestor de Plantillas: 14,700+ líneas que definen estructura de reportes",
    "Sistema de borradores con flujo: editar → enviar → revisar → aprobar → COI",
    "Gráficas interactivas con Chart.js + snapshots en localStorage",
    "Seguridad: bcrypt(12), JWT, rate limiting, helmet, CORS, Joi, SQL parametrizado",
    "Backups automáticos cada 60 min con SHA256 + auto-update vía GitHub Releases",
    "",
    "Versión 4.1.0  |  Febrero 2026  |  Iustus Renidet & J02V4N",
], 11, AZUL_CLARO)


# === GUARDAR ===
output_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                           'PRESENTACION_TECNICA_DEFINITIVA_SUMMACHAM.pptx')
prs.save(output_path)
print(f"Presentacion guardada: {output_path}")
print(f"Total slides: {len(prs.slides)}")
